from __future__ import annotations

from pathlib import Path

from ..models import AttachmentReaderConfig


class OfficeAttachmentReader:
    """
    負責在 builder.attachment.readers.office_reader 中封裝 OfficeAttachmentReader，封裝附件讀取與內容萃取流程，將檔案轉成可推理的證據。
    
    Args:
        config: 控制此流程行為的設定資料。
    
    Returns:
        類別本身不直接回傳值；建立實例後可透過其方法操作狀態與流程。
    
    限制或副作用:
        方法可能更新內部狀態、讀寫檔案、呼叫外部服務或產生日誌，需依使用情境確認。
    """
    def __init__(self, config: AttachmentReaderConfig) -> None:
        """
        負責執行 OfficeAttachmentReader 中的 __init__ 流程，初始化物件所需的設定、依賴與內部狀態，讓後續方法可以沿用同一份執行上下文。
        
        Args:
            config: 控制此流程行為的設定資料。
        
        Returns:
            執行結果；若函式標註回傳型別，預期型別為 None。
        
        限制或副作用:
            可能讀取或更新物件狀態、檔案、外部服務或日誌；請依呼叫場景確認副作用。
        """
        self.config = config

    def read_docx(self, file_path: Path) -> str:
        """
        負責執行 OfficeAttachmentReader 中的 read_docx 流程，讀取本地或外部資料來源並轉換成系統可處理的格式。
        
        Args:
            file_path: 要讀取或寫入的檔案或目錄路徑。
        
        Returns:
            執行結果；若函式標註回傳型別，預期型別為 str。
        
        限制或副作用:
            可能讀取或更新物件狀態、檔案、外部服務或日誌；請依呼叫場景確認副作用。
        """
        try:
            from docx import Document
        except ImportError:
            return "DOCX attachment detected, but python-docx is not installed."

        doc = Document(str(file_path))
        lines: list[str] = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                lines.append(text)
        for table_idx, table in enumerate(doc.tables, start=1):
            lines.append(f"[table {table_idx}]")
            for row in table.rows[: self.config.max_table_rows]:
                lines.append(" | ".join(cell.text.strip() for cell in row.cells))
        return "\n".join(lines)

    def read_pptx(self, file_path: Path) -> str:
        """
        負責執行 OfficeAttachmentReader 中的 read_pptx 流程，讀取本地或外部資料來源並轉換成系統可處理的格式。
        
        Args:
            file_path: 要讀取或寫入的檔案或目錄路徑。
        
        Returns:
            執行結果；若函式標註回傳型別，預期型別為 str。
        
        限制或副作用:
            可能讀取或更新物件狀態、檔案、外部服務或日誌；請依呼叫場景確認副作用。
        """
        try:
            from pptx import Presentation
        except ImportError:
            return "PPTX attachment detected, but python-pptx is not installed."

        presentation = Presentation(str(file_path))
        lines: list[str] = []
        for slide_idx, slide in enumerate(presentation.slides, start=1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    slide_lines.append(" ".join(text.split()))
            if slide_lines:
                lines.append(f"[slide {slide_idx}] " + "\n".join(slide_lines))
        return "\n".join(lines)

    def read_pdf(self, file_path: Path) -> str:
        """
        負責執行 OfficeAttachmentReader 中的 read_pdf 流程，讀取本地或外部資料來源並轉換成系統可處理的格式。
        
        Args:
            file_path: 要讀取或寫入的檔案或目錄路徑。
        
        Returns:
            執行結果；若函式標註回傳型別，預期型別為 str。
        
        限制或副作用:
            可能讀取或更新物件狀態、檔案、外部服務或日誌；請依呼叫場景確認副作用。
        """
        reader_cls = None
        try:
            from pypdf import PdfReader

            reader_cls = PdfReader
        except ImportError:
            try:
                from PyPDF2 import PdfReader

                reader_cls = PdfReader
            except ImportError:
                return "PDF attachment detected, but pypdf/PyPDF2 is not installed."

        reader = reader_cls(str(file_path))
        lines: list[str] = []
        for idx, page in enumerate(reader.pages[: self.config.max_pdf_pages], start=1):
            text = page.extract_text() or ""
            if text.strip():
                lines.append(f"[page {idx}]\n{text.strip()}")
        if len(reader.pages) > self.config.max_pdf_pages:
            lines.append(f"[truncated after {self.config.max_pdf_pages} pages]")
        return "\n\n".join(lines)
